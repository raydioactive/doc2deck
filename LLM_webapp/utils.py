# LLM_webapp/utils.py

import docx2txt
import genanki
import random
import csv # Needed for TSV/CSV export
import io # Needed for TSV/CSV export in memory
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
import markdown 
import re


def extract_text_from_pdf(pdf_path: str) -> Optional[str]:
    """Extract text from PDF using pypdf with improved robustness."""
    try:
        with open(pdf_path, 'rb') as file:
            reader = PdfReader(file)
            text_parts = []
            
            # Process each page
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    # Basic cleaning
                    page_text = re.sub(r'\s+', ' ', page_text)  # Normalize whitespace
                    page_text = re.sub(r'([a-z])-\s+([a-z])', r'\1\2', page_text)  # Fix hyphenation
                    text_parts.append(page_text)
            
            full_text = "\n\n".join(text_parts).strip()
            return full_text if full_text else None
            
    except FileNotFoundError:
        # Re-raise FileNotFoundError to be handled by the caller
        raise
    except Exception as e:
        # Print an error message and return None for other pypdf errors
        print(f"Extraction failed using pypdf for {pdf_path}: {e}")
        return None

def extract_text_from_docx(file_path: str) -> str:
    """Extract text content from a DOCX file"""
    try:
        return docx2txt.process(file_path)
    except Exception as e:
        raise IOError(f"Failed to extract text from DOCX file: {file_path}. Error: {e}") from e

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text content from a PPTX file with improved table and chart handling"""
    try:
        prs = Presentation(file_path)
        text_runs = []
        
        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            # Add slide number reference
            text_runs.append(f"--- Slide {slide_num} ---")
            
            # Extract slide title if available
            if slide.shapes.title:
                text_runs.append(f"Title: {slide.shapes.title.text}")
            
            # Process all shapes
            for shape in slide.shapes:
                # Text frames
                if shape.has_text_frame:
                    shape_text = []
                    for paragraph in shape.text_frame.paragraphs:
                        shape_text.append(paragraph.text)
                    if shape_text:
                        text_runs.append("\n".join(shape_text))
                
                # Tables with improved extraction
                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        row_cells = []
                        for cell in row.cells:
                            if cell.text_frame:
                                row_cells.append(cell.text_frame.text.strip())
                            else:
                                row_cells.append("")
                        table_rows.append(" | ".join(row_cells))
                    if table_rows:
                        text_runs.append("Table:\n" + "\n".join(table_rows))
                
                # Charts with more detailed extraction
                if hasattr(shape, 'chart'):
                    try:
                        chart_text = ["Chart:"]
                        
                        # Get chart title
                        if shape.chart.has_title:
                            chart_text.append(f"Title: {shape.chart.chart_title.text_frame.text}")
                        
                        # Extract category labels if possible
                        if hasattr(shape.chart, 'plots') and shape.chart.plots:
                            for plot in shape.chart.plots:
                                if hasattr(plot, 'categories') and plot.categories:
                                    categories = []
                                    for cat in plot.categories:
                                        if cat:
                                            categories.append(str(cat))
                                    if categories:
                                        chart_text.append(f"Categories: {', '.join(categories)}")
                        
                        text_runs.append("\n".join(chart_text))
                    except Exception as e:
                        text_runs.append(f"Chart: [Data extraction error: {str(e)}]")

        return "\n\n".join(text_runs).strip() # Join text with newlines
    except Exception as e:
        raise IOError(f"Failed to extract text from PPTX file: {file_path}. Error: {e}") from e

# --- Anki Package (.apkg) Exporter ---
def save_flashcards_as_apkg(flashcards: List[Dict[str, str]], deck_name: str, output_path: str) -> None:
    """
    Generates an Anki .apkg file from a list of flashcards, converting Markdown lists to HTML.
    """
    if not deck_name:
        deck_name = "Generated Deck" # Default deck name if none provided

    # Generate unique IDs for model and deck
    model_id = random.randrange(1 << 30, 1 << 31)
    deck_id = random.randrange(1 << 30, 1 << 31)

    # Define the Anki model (defines card structure and appearance)
    anki_model = genanki.Model(
        model_id,
        f'Simple Model (Genanki {model_id})',
        fields=[
            {'name': 'Front'},
            {'name': 'Back'},
        ],
        templates=[
            {
                'name': 'Card 1',
                'qfmt': '{{Front}}', # Question format - Now expects HTML
                'afmt': '{{FrontSide}}\n\n<hr id="answer">\n\n{{Back}}', # Answer format - Now expects HTML
            },
        ],
        # Basic CSS styling (optional) - Anki's default list styling might apply
        css='''
            .card { font-family: arial; font-size: 20px; text-align: center; color: black; background-color: white; }
            hr { border-top: 1px solid #ccc; }
            /* You might add specific list styling here if needed */
            ul, ol { text-align: left; display: inline-block; } /* Example */
            '''
    )

    # Create the Anki deck
    anki_deck = genanki.Deck(deck_id, deck_name)

    # Create Anki notes from the flashcards
    for card in flashcards:
        front_text = card.get('front', '').strip()
        back_text = card.get('back', '').strip()

        # Convert potential Markdown in front and back fields to HTML
        # Using 'nl2br' extension converts newlines to <br> tags, which is often desired in Anki
        # Using 'fenced_code' allows for code blocks if the LLM generates them
        front_html = markdown.markdown(front_text, extensions=['nl2br', 'fenced_code']) if front_text else ''
        back_html = markdown.markdown(back_text, extensions=['nl2br', 'fenced_code']) if back_text else ''

        # Only add cards that have some content on the front (after potential HTML conversion)
        if front_html: # Check the HTML version
            anki_note = genanki.Note(
                model=anki_model,
                fields=[front_html, back_html] # Use the HTML versions
            )
            anki_deck.add_note(anki_note)

    # Create the package and save to file
    try:
        anki_package = genanki.Package(anki_deck)
        anki_package.write_to_file(output_path)
    except Exception as e:
        raise IOError(f"Failed to write Anki deck to '{output_path}'. Error: {e}") from e
    
# --- Text File (.txt / TSV) Exporter ---
def save_flashcards_as_tsv(flashcards: List[Dict[str, str]], output_path: str) -> None:
    """
    Saves flashcards as a Tab-Separated Value (TSV) text file.
    Each line represents a card: Front<TAB>Back.
    """
    try:
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f, delimiter='\t', lineterminator='\n')
            for card in flashcards:
                front = card.get('front', '').strip()
                back = card.get('back', '').strip()
                # Replace newlines within fields with spaces or HTML breaks for Anki import
                front = front.replace('\n', ' ')
                back = back.replace('\n', '<br>') # Anki understands <br> in TSV
                writer.writerow([front, back])
    except Exception as e:
        raise IOError(f"Failed to write TSV file to '{output_path}'. Error: {e}") from e

def generate_tsv_content(flashcards: List[Dict[str, str]]) -> str:
    """
    Generates TSV content as a string (useful for direct download).
    """
    output = io.StringIO()
    writer = csv.writer(output, delimiter='\t', lineterminator='\n')
    for card in flashcards:
        front = card.get('front', '').strip().replace('\n', ' ')
        back = card.get('back', '').strip().replace('\n', '<br>')
        writer.writerow([front, back])
    return output.getvalue()

# --- CSV Exporter Functions ---
def save_flashcards_as_csv(flashcards: List[Dict[str, str]], output_path: str) -> None:
    """
    Saves flashcards as a Comma Separated Value (CSV) file.
    Includes a header row: Front,Back.
    Standard CSV quoting applied automatically for fields containing commas or newlines.
    """
    try:
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            # Use default comma delimiter
            writer = csv.writer(f, delimiter=',', quotechar='"', quoting=csv.QUOTE_MINIMAL)
            # Write header
            writer.writerow(['Front', 'Back'])
            # Write card data
            for card in flashcards:
                front = card.get('front', '').strip()
                back = card.get('back', '').strip()
                writer.writerow([front, back])
    except Exception as e:
        raise IOError(f"Failed to write CSV file to '{output_path}'. Error: {e}") from e

def generate_csv_content(flashcards: List[Dict[str, str]]) -> str:
    """
    Generates CSV content as a string, including a header row.
    """
    output = io.StringIO()
    writer = csv.writer(output, delimiter=',', quotechar='"', quoting=csv.QUOTE_MINIMAL)
    writer.writerow(['Front', 'Back']) # Header row
    for card in flashcards:
        front = card.get('front', '').strip()
        back = card.get('back', '').strip()
        writer.writerow([front, back])
    return output.getvalue()

def filter_low_quality_cards(cards: List[Dict[str, str]]) -> List[Dict[str, str]]:
    """Filter out low-quality flashcards based on various criteria."""
    filtered_cards = []
    
    for card in cards:
        front = card.get('front', '').strip()
        back = card.get('back', '').strip()
        
        # Skip empty cards
        if not front or not back:
            continue
            
        # Skip cards with very short answers
        if len(back) < 10:
            continue
            
        # Skip cards with very similar front/back content
        if similarity_ratio(front, back) > 0.8:
            continue
            
        # Skip cards with certain problematic patterns
        if is_problematic_card(front, back):
            continue
            
        filtered_cards.append(card)
        
    return filtered_cards

def similarity_ratio(str1: str, str2: str) -> float:
    """Calculate similarity between two strings using sequence matching."""
    from difflib import SequenceMatcher
    return SequenceMatcher(None, str1.lower(), str2.lower()).ratio()

def is_problematic_card(front: str, back: str) -> bool:
    """Check for common problems in flashcards."""
    # Front contains the answer (defeats the purpose)
    if back.strip() in front.strip():
        return True
        
    # Back just repeats the front with minor changes
    if front.lower().strip() == back.lower().strip():
        return True
        
    # Front is a complete sentence and back is just True/False
    if len(front) > 30 and back.strip().lower() in ['true', 'false']:
        return True
        
    # Card is just a definition format with no actual content
    definition_patterns = [
        r"^Definition of ",
        r"^Define ",
        r"^What is the definition of "
    ]
    
    for pattern in definition_patterns:
        if re.match(pattern, front, re.IGNORECASE) and len(back) < 20:
            return True
            
    return False